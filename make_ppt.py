from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
PRIMARY = RGBColor(0x1B, 0x3A, 0x5C)      # Dark navy
ACCENT = RGBColor(0x2E, 0x86, 0xC1)       # Blue
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2C, 0x3E, 0x50)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
ACCENT2 = RGBColor(0xE7, 0x4C, 0x3C)      # Red accent
ACCENT3 = RGBColor(0x27, 0xAE, 0x60)      # Green
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text="", font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_text(text_frame, text, font_size=16, color=DARK, bold=False, level=0, space_before=Pt(6), font_name="Microsoft YaHei"):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.level = level
    p.space_before = space_before
    return p

# ============================================================
# SLIDE 1: Cover
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Background
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), PRIMARY)

# Decorative accent bar
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.0), Inches(13.333), Inches(0.06), ACCENT)

# Title
add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
            "月度工作进展汇报", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle line
add_textbox(slide, Inches(2), Inches(3.3), Inches(9), Inches(0.8),
            "基于多元时序曲线相似性的分布式集群异常检测方法研究", font_size=20, color=RGBColor(0xAE, 0xBF, 0xD5), alignment=PP_ALIGN.CENTER)

# Info
add_textbox(slide, Inches(3), Inches(4.8), Inches(7), Inches(0.5),
            "汇报人：XXX    |    导师：XXX", font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC), alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(3), Inches(5.4), Inches(7), Inches(0.5),
            "中国科学院大学  /  中国联通（联合培养）", font_size=16, color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: 大论文
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Header bar
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY)
add_textbox(slide, Inches(0.6), Inches(0.2), Inches(10), Inches(0.7),
            "01  硕士学位论文 -- 撰写修改与盲审提交", font_size=30, color=WHITE, bold=True)

# Left column - 论文信息
left_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5), LIGHT_BG)
left_box.adjustments[0] = 0.02

tb = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.3), Inches(0.5),
                 "论文信息", font_size=20, color=PRIMARY, bold=True)

tb2 = add_textbox(slide, Inches(0.8), Inches(2.2), Inches(5.3), Inches(4.5), "", font_size=15, color=DARK)
tf = tb2.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "题目：基于多元时序曲线相似性的分布式集群"
tf.paragraphs[0].font.size = Pt(15)
tf.paragraphs[0].font.name = "Microsoft YaHei"
tf.paragraphs[0].font.color.rgb = DARK
add_bullet_text(tf, "          异常检测方法研究", font_size=15, color=DARK)
add_bullet_text(tf, "", font_size=8)
add_bullet_text(tf, "领域：软件工程 / 智能运维（AIOps）", font_size=15, color=DARK)
add_bullet_text(tf, "", font_size=8)
add_bullet_text(tf, "类型：方法研究 + 系统实现", font_size=15, color=DARK)
add_bullet_text(tf, "", font_size=8)
add_bullet_text(tf, "结构：共六章", font_size=15, color=DARK)
add_bullet_text(tf, "  - 绪论", font_size=14, color=GRAY, level=1)
add_bullet_text(tf, "  - 背景知识与理论基础", font_size=14, color=GRAY, level=1)
add_bullet_text(tf, "  - 单指标多节点异常检测", font_size=14, color=GRAY, level=1)
add_bullet_text(tf, "  - 多指标多节点异常检测", font_size=14, color=GRAY, level=1)
add_bullet_text(tf, "  - 系统实现", font_size=14, color=GRAY, level=1)
add_bullet_text(tf, "  - 总结与展望", font_size=14, color=GRAY, level=1)

# Right column - 本月工作
right_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.5), LIGHT_BG)
right_box.adjustments[0] = 0.02

add_textbox(slide, Inches(7.1), Inches(1.6), Inches(5.5), Inches(0.5),
            "本月完成工作", font_size=20, color=PRIMARY, bold=True)

tb3 = add_textbox(slide, Inches(7.1), Inches(2.2), Inches(5.5), Inches(4.5), "", font_size=15, color=DARK)
tf3 = tb3.text_frame
tf3.word_wrap = True
tf3.paragraphs[0].text = "✦  完成全文六章内容的撰写与反复修改"
tf3.paragraphs[0].font.size = Pt(15)
tf3.paragraphs[0].font.name = "Microsoft YaHei"
tf3.paragraphs[0].font.color.rgb = DARK
add_bullet_text(tf3, "", font_size=6)
add_bullet_text(tf3, "✦  重点修改内容：", font_size=15, color=DARK, bold=True)
add_bullet_text(tf3, "    - 补充有效性威胁讨论", font_size=14, color=DARK, level=1)
add_bullet_text(tf3, "    - 完善实验对比与消融实验", font_size=14, color=DARK, level=1)
add_bullet_text(tf3, "    - 优化图表与系统架构可视化", font_size=14, color=DARK, level=1)
add_bullet_text(tf3, "", font_size=6)
add_bullet_text(tf3, "✦  完成格式审查，符合学位论文模板规范", font_size=15, color=DARK)
add_bullet_text(tf3, "", font_size=6)
add_bullet_text(tf3, "✦  已按时提交盲审", font_size=15, color=DARK, bold=True)

# Status badge
status = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(5.8), Inches(2.5), Inches(0.55), ACCENT3)
status.adjustments[0] = 0.3
add_textbox(slide, Inches(9.0), Inches(5.82), Inches(2.5), Inches(0.5),
            "当前状态：盲审中", font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 3: 小论文
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY)
add_textbox(slide, Inches(0.6), Inches(0.2), Inches(10), Inches(0.7),
            "02  学术小论文 -- 初稿撰写", font_size=30, color=WHITE, bold=True)

# Left - 论文框架
left_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5), LIGHT_BG)
left_box.adjustments[0] = 0.02

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.3), Inches(0.5),
            "论文方向与框架", font_size=20, color=PRIMARY, bold=True)

tb = add_textbox(slide, Inches(0.8), Inches(2.2), Inches(5.3), Inches(4.5), "", font_size=15, color=DARK)
tf = tb.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "方向：基于分布式集群多节点曲线相似性的"
tf.paragraphs[0].font.size = Pt(15)
tf.paragraphs[0].font.name = "Microsoft YaHei"
tf.paragraphs[0].font.color.rgb = DARK
add_bullet_text(tf, "          异常检测方法", font_size=15)
add_bullet_text(tf, "", font_size=6)
add_bullet_text(tf, "来源：从大论文核心章节提炼", font_size=15, color=GRAY)
add_bullet_text(tf, "", font_size=10)
add_bullet_text(tf, "论文结构：", font_size=16, bold=True, color=PRIMARY)
add_bullet_text(tf, "", font_size=6)
add_bullet_text(tf, "  1. Introduction", font_size=15)
add_bullet_text(tf, "  2. Related Work", font_size=15)
add_bullet_text(tf, "  3. Method", font_size=15)
add_bullet_text(tf, "  4. Experiments", font_size=15)
add_bullet_text(tf, "  5. Conclusion", font_size=15)

# Right - 核心内容与计划
right_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.5), LIGHT_BG)
right_box.adjustments[0] = 0.02

add_textbox(slide, Inches(7.1), Inches(1.6), Inches(5.5), Inches(0.5),
            "核心内容与后续计划", font_size=20, color=PRIMARY, bold=True)

tb = add_textbox(slide, Inches(7.1), Inches(2.2), Inches(5.5), Inches(4.5), "", font_size=15, color=DARK)
tf = tb.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "已完成的核心内容："
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.name = "Microsoft YaHei"
tf.paragraphs[0].font.color.rgb = PRIMARY
tf.paragraphs[0].font.bold = True
add_bullet_text(tf, "", font_size=4)
add_bullet_text(tf, "✦  三维分析范式的创新点凝练", font_size=15)
add_bullet_text(tf, "     （时间 x 指标 x 节点）", font_size=13, color=GRAY)
add_bullet_text(tf, "✦  曲线相似性异常检测算法描述", font_size=15)
add_bullet_text(tf, "✦  真实数据集实验结果（F1 > 0.97）", font_size=15, bold=True, color=ACCENT2)
add_bullet_text(tf, "", font_size=10)
add_bullet_text(tf, "后续计划：", font_size=16, bold=True, color=PRIMARY)
add_bullet_text(tf, "", font_size=4)
add_bullet_text(tf, "✦  根据导师意见修改完善", font_size=15)
add_bullet_text(tf, "✦  补充与最新相关工作的对比讨论", font_size=15)
add_bullet_text(tf, "✦  目标：下月完成投稿", font_size=15, bold=True)

# ============================================================
# SLIDE 4: 企业实践
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY)
add_textbox(slide, Inches(0.6), Inches(0.2), Inches(10), Inches(0.7),
            "03  企业实践 -- 联通智能运维工作", font_size=30, color=WHITE, bold=True)

# Left - 异常分析复盘
left_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(3.8), Inches(5.5), LIGHT_BG)
left_box.adjustments[0] = 0.02

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(3.3), Inches(0.5),
            "异常分析复盘", font_size=20, color=PRIMARY, bold=True)

tb = add_textbox(slide, Inches(0.8), Inches(2.3), Inches(3.3), Inches(4.2), "", font_size=15, color=DARK)
tf = tb.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "✦  参与生产环境真实异常"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.name = "Microsoft YaHei"
tf.paragraphs[0].font.color.rgb = DARK
add_bullet_text(tf, "    事件的回顾分析", font_size=14)
add_bullet_text(tf, "", font_size=6)
add_bullet_text(tf, "✦  对历史告警案例进行", font_size=14)
add_bullet_text(tf, "    根因归类与总结", font_size=14)
add_bullet_text(tf, "", font_size=6)
add_bullet_text(tf, "✦  积累运维场景经验，", font_size=14)
add_bullet_text(tf, "    为算法优化提供反馈", font_size=14)

# Middle - 研发项目验证
mid_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.7), Inches(1.5), Inches(3.9), Inches(5.5), LIGHT_BG)
mid_box.adjustments[0] = 0.02

add_textbox(slide, Inches(5.0), Inches(1.6), Inches(3.5), Inches(0.5),
            "研发项目后续验证", font_size=20, color=PRIMARY, bold=True)

tb = add_textbox(slide, Inches(5.0), Inches(2.3), Inches(3.5), Inches(4.2), "", font_size=15, color=DARK)
tf = tb.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "✦  对已部署异常检测算法"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.name = "Microsoft YaHei"
tf.paragraphs[0].font.color.rgb = DARK
add_bullet_text(tf, "    进行效果跟踪评估", font_size=14)
add_bullet_text(tf, "", font_size=6)
add_bullet_text(tf, "✦  收集新场景验证数据，", font_size=14)
add_bullet_text(tf, "    评估模型泛化性", font_size=14)
add_bullet_text(tf, "", font_size=6)
add_bullet_text(tf, "✦  配合团队完成检测结果", font_size=14)
add_bullet_text(tf, "    的精度与召回分析", font_size=14)

# Right - 实践收获
right_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(1.5), Inches(3.8), Inches(5.5), ACCENT)
right_box.adjustments[0] = 0.02

add_textbox(slide, Inches(9.3), Inches(1.6), Inches(3.3), Inches(0.5),
            "实践收获", font_size=20, color=WHITE, bold=True)

tb = add_textbox(slide, Inches(9.3), Inches(2.3), Inches(3.3), Inches(4.2), "", font_size=15, color=WHITE)
tf = tb.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "加深了对分布式集群真实"
tf.paragraphs[0].font.size = Pt(15)
tf.paragraphs[0].font.name = "Microsoft YaHei"
tf.paragraphs[0].font.color.rgb = WHITE
add_bullet_text(tf, "运维场景的理解", font_size=15, color=WHITE)
add_bullet_text(tf, "", font_size=12)
add_bullet_text(tf, "论文方法与实际业务需求", font_size=15, color=WHITE)
add_bullet_text(tf, "形成了良好的闭环验证", font_size=15, color=WHITE)

# ============================================================
# SLIDE 5: 下月计划
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY)
add_textbox(slide, Inches(0.6), Inches(0.2), Inches(10), Inches(0.7),
            "04  下月工作计划", font_size=30, color=WHITE, bold=True)

items = [
    ("大论文", "关注盲审进度，准备根据评审意见进行修改", ACCENT),
    ("小论文", "完成修改并争取投稿", ACCENT3),
    ("企业实践", "继续参与智能运维项目，推进检测算法在新场景的验证", RGBColor(0xF3, 0x9C, 0x12)),
    ("答辩准备", "着手准备学位答辩材料与PPT", ACCENT2),
]

for i, (title, desc, color) in enumerate(items):
    y = Inches(1.6) + Inches(1.35) * i
    # Color bar
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), y, Inches(0.15), Inches(1.0), color)
    # Card background
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), y, Inches(10.5), Inches(1.0), LIGHT_BG)
    card.adjustments[0] = 0.02
    # Number
    add_textbox(slide, Inches(1.6), y + Inches(0.05), Inches(0.8), Inches(0.9),
                str(i+1), font_size=32, color=color, bold=True)
    # Title
    add_textbox(slide, Inches(2.3), y + Inches(0.05), Inches(2.0), Inches(0.45),
                title, font_size=20, color=PRIMARY, bold=True)
    # Description
    add_textbox(slide, Inches(2.3), y + Inches(0.5), Inches(9.3), Inches(0.45),
                desc, font_size=16, color=GRAY)

# ============================================================
# SLIDE 6: 致谢
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), PRIMARY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.4), Inches(13.333), Inches(0.06), ACCENT)

add_textbox(slide, Inches(2), Inches(2.2), Inches(9), Inches(1.2),
            "感谢各位老师指导！", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(4.0), Inches(9), Inches(0.8),
            "欢迎批评指正", font_size=24, color=RGBColor(0xAE, 0xBF, 0xD5), alignment=PP_ALIGN.CENTER)

# Save
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "月度工作汇报.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
